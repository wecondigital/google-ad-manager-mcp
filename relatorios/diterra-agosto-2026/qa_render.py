#!/usr/bin/env python3
"""Renderizador de QA: le o .pptx gerado e desenha cada slide com PIL.

Existe porque o LibreOffice esta quebrado neste container. Nao substitui um
render fiel do PowerPoint, mas expoe o que interessa no QA de layout:
posicao, tamanho, sobreposicao e estouro de texto.

Usa Liberation Sans/Serif (mais largas que Calibri/Cambria), entao a checagem
de estouro e conservadora: o que couber aqui, cabe no PowerPoint.
"""
import sys, glob, os
from pptx import Presentation
from pptx.util import Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image, ImageDraw, ImageFont

SCALE = 105.0  # px por polegada
FONTS = {
    (False, False): "/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf",
    (True,  False): "/usr/share/fonts/truetype/liberation/LiberationSans-Bold.ttf",
    (False, True):  "/usr/share/fonts/truetype/liberation/LiberationSerif-Regular.ttf",
    (True,  True):  "/usr/share/fonts/truetype/liberation/LiberationSerif-Bold.ttf",
}
_cache = {}


def font(size_pt, bold, serif):
    key = (round(size_pt * 2), bold, serif)
    if key not in _cache:
        _cache[key] = ImageFont.truetype(FONTS[(bold, serif)], max(6, int(round(size_pt * SCALE / 72.0))))
    return _cache[key]


def emu_px(v):
    return None if v is None else Emu(v).inches * SCALE


def rgb(color, default=(30, 34, 42)):
    try:
        if color and color.type is not None and color.rgb is not None:
            return tuple(bytes.fromhex(str(color.rgb)))
    except Exception:
        pass
    return default


def shape_fill(sh):
    try:
        f = sh.fill
        if f.type is not None and f.type == 1:  # solid
            return rgb(f.fore_color, None)
    except Exception:
        pass
    return None


def shape_line(sh):
    try:
        ln = sh.line
        if ln.fill.type == 1:
            return rgb(ln.color, None)
    except Exception:
        pass
    return None


def wrap(draw, text, fnt, max_w):
    lines = []
    for para in text.split("\n"):
        if not para:
            lines.append("")
            continue
        words, cur = para.split(" "), ""
        for w in words:
            trial = (cur + " " + w).strip()
            if draw.textlength(trial, font=fnt) <= max_w or not cur:
                cur = trial
            else:
                lines.append(cur)
                cur = w
        lines.append(cur)
    return lines


def render(pptx_path, outdir):
    prs = Presentation(pptx_path)
    SW = Emu(prs.slide_width).inches * SCALE
    SH = Emu(prs.slide_height).inches * SCALE
    problems = []

    for idx, slide in enumerate(prs.slides, 1):
        bg = (255, 255, 255)
        # fundo do slide vem do <p:bg> — python-pptx nao expoe, lemos o XML
        xml = slide._element.xml
        if "<p:bg>" in xml:
            import re
            m = re.search(r"<p:bg>.*?srgbClr val=\"([0-9A-Fa-f]{6})\"", xml, re.S)
            if m:
                bg = tuple(bytes.fromhex(m.group(1)))
        img = Image.new("RGB", (int(SW), int(SH)), bg)
        d = ImageDraw.Draw(img)

        for sh in slide.shapes:
            x, y = emu_px(sh.left), emu_px(sh.top)
            w, h = emu_px(sh.width), emu_px(sh.height)
            if x is None:
                continue
            box = (x, y, x + w, y + h)

            if x < -1 or y < -1 or x + w > SW + 1 or y + h > SH + 1:
                problems.append(f"slide {idx}: '{(sh.name or '')[:34]}' fora dos limites "
                                f"({x/SCALE:.2f},{y/SCALE:.2f} .. {(x+w)/SCALE:.2f},{(y+h)/SCALE:.2f})")

            if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    pic = Image.open(io_bytes(sh)).convert("RGBA")
                    pic = pic.resize((max(1, int(w)), max(1, int(h))))
                    img.paste(pic, (int(x), int(y)), pic)
                except Exception:
                    d.rectangle(box, outline=(150, 150, 150))
                continue

            if getattr(sh, "has_table", False) and sh.has_table:
                tbl = sh.table
                cols = [Emu(c.width).inches * SCALE for c in tbl.columns]
                rows = [Emu(r.height).inches * SCALE for r in tbl.rows]
                cy = y
                for ri, row in enumerate(tbl.rows):
                    cx = x
                    for ci, cell in enumerate(row.cells):
                        cw, ch = cols[ci], rows[ri]
                        cf = None
                        try:
                            if cell.fill.type == 1:
                                cf = rgb(cell.fill.fore_color, None)
                        except Exception:
                            pass
                        d.rectangle((cx, cy, cx + cw, cy + ch), fill=cf, outline=(218, 215, 198))
                        runs = [r for p in cell.text_frame.paragraphs for r in p.runs]
                        if runs:
                            rf = runs[0].font
                            sz = rf.size.pt if rf.size else 11
                            fnt2 = font(sz, bool(rf.bold), bool(rf.name and "Cambria" in rf.name))
                            cc = rgb(rf.color, (30, 34, 42))
                            txt = "".join(r.text for r in runs)
                            al = str(cell.text_frame.paragraphs[0].alignment or "LEFT")
                            tw = d.textlength(txt, font=fnt2)
                            if "CENTER" in al:
                                tx = cx + (cw - tw) / 2
                            elif "RIGHT" in al:
                                tx = cx + cw - tw - 8
                            else:
                                tx = cx + 8
                            d.text((tx, cy + (ch - sz * SCALE / 72.0 * 1.2) / 2), txt, font=fnt2, fill=cc)
                            if tw > cw - 12:
                                problems.append(
                                    f"slide {idx}: célula da tabela estoura (col {ci}, "
                                    f"{tw/SCALE:.2f}\" em {cw/SCALE:.2f}\") — \"{txt}\"")
                        cx += cw
                    cy += rows[ri]
                continue

            if sh.has_chart or sh.shape_type == MSO_SHAPE_TYPE.CHART:
                d.rectangle(box, fill=(246, 246, 244), outline=(190, 190, 185))
                d.text((x + 8, y + 8), "[gráfico nativo]", font=font(11, True, False), fill=(120, 120, 118))
                continue

            fill, line = shape_fill(sh), shape_line(sh)
            if fill or line:
                d.rectangle(box, fill=fill, outline=line)

            if not sh.has_text_frame:
                continue
            tf = sh.text_frame
            runs = [(r.text, r.font) for p in tf.paragraphs for r in p.runs]
            if not all(t for t, _ in runs) and not any(t.strip() for t, _ in runs):
                continue
            text = "\n".join("".join(r.text for r in p.runs) for p in tf.paragraphs)
            if not text.strip():
                continue

            f0 = runs[0][1] if runs else None
            size = f0.size.pt if (f0 and f0.size) else 12
            bold = bool(f0 and f0.bold)
            serif = bool(f0 and f0.name and "Cambria" in (f0.name or ""))
            col = rgb(f0.color, (30, 34, 42)) if f0 else (30, 34, 42)
            fnt = font(size, bold, serif)

            pad = 2
            lines = wrap(d, text, fnt, max(4, w - 2 * pad))
            lh = size * SCALE / 72.0 * 1.20
            th = lh * len(lines)
            align = (tf.paragraphs[0].alignment.__str__() if tf.paragraphs[0].alignment else "LEFT")
            valign = str(tf.vertical_anchor) if tf.vertical_anchor else "TOP"

            ty = y + (h - th) / 2 if "MIDDLE" in valign else (y + h - th - pad if "BOTTOM" in valign else y + pad)
            for ln in lines:
                lw = d.textlength(ln, font=fnt)
                tx = x + (w - lw) / 2 if "CENTER" in align else (x + w - lw - pad if "RIGHT" in align else x + pad)
                d.text((tx, ty), ln, font=fnt, fill=col)
                ty += lh

            if th > h + 1.5:
                problems.append(f"slide {idx}: texto estoura a caixa ({th/SCALE:.2f}\" em "
                                f"{h/SCALE:.2f}\") — \"{text[:52]}…\"")
            longest = max((d.textlength(l, font=fnt) for l in lines), default=0)
            if longest > w - 2 * pad + 1.5:
                problems.append(f"slide {idx}: linha mais larga que a caixa — \"{text[:52]}…\"")

        img.save(os.path.join(outdir, f"slide-{idx}.png"))

    print(f"{len(prs.slides.__iter__.__self__._sldIdLst)} slides renderizados em {outdir}")
    if problems:
        print("\n=== PROBLEMAS ===")
        for p in problems:
            print(" -", p)
    else:
        print("\nNenhum estouro ou shape fora dos limites.")


def io_bytes(sh):
    import io
    return io.BytesIO(sh.image.blob)


if __name__ == "__main__":
    render(sys.argv[1], sys.argv[2] if len(sys.argv) > 2 else ".")
