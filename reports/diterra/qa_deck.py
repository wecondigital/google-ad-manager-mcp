# -*- coding: utf-8 -*-
"""QA geometrico do deck: estouro de texto, sobreposicao e margem.

LibreOffice nao roda neste ambiente, entao em vez de olhar o render eu meco a
geometria: largura do texto estimada por metrica da Arial contra a caixa que o
contem, e retangulos de conteudo cruzados dois a dois.
"""
import glob
import os
import sys

from pptx import Presentation
from pptx.util import Emu

DECK = sys.argv[1] if len(sys.argv) > 1 else sorted(
    glob.glob(os.path.join(os.path.dirname(os.path.abspath(__file__)),
                           "Acompanhamento-DiTerra-*.pptx")))[-1]
EMU = 914400.0
LARG, ALT = 10.0, 5.625
MARGEM = 0.5

# largura media de caractere na Arial, em fracao do tamanho da fonte
K_REG, K_BOLD = 0.52, 0.56
ALTURA_LINHA = 1.22   # entrelinha aproximada


def caixas(slide):
    out = []
    for sh in slide.shapes:
        if sh.left is None or sh.top is None:
            continue
        out.append({
            "sh": sh,
            "x": sh.left / EMU, "y": sh.top / EMU,
            "w": (sh.width or 0) / EMU, "h": (sh.height or 0) / EMU,
            "tipo": sh.shape_type,
            "texto": (sh.text_frame.text if sh.has_text_frame else ""),
        })
    return out


def altura_estimada(c):
    """Quantas linhas o texto ocupa e qual altura isso pede."""
    sh = c["sh"]
    if not sh.has_text_frame or not c["texto"].strip():
        return 0.0, 0
    linhas = 0
    altura = 0.0
    for par in sh.text_frame.paragraphs:
        txt = "".join(r.text for r in par.runs)
        if not txt:
            linhas += 1
            continue
        pt = max([(r.font.size.pt if r.font.size else 12) for r in par.runs] or [12])
        bold = any(r.font.bold for r in par.runs)
        k = K_BOLD if bold else K_REG
        larg_char = pt * k / 72.0
        cabe = max(1, int(c["w"] / larg_char)) if larg_char else 999
        n = max(1, -(-len(txt) // cabe))          # ceil
        linhas += n
        altura += n * pt * ALTURA_LINHA / 72.0
    return altura, linhas


def sobrepoe(a, b):
    return not (a["x"] + a["w"] <= b["x"] + 1e-6 or b["x"] + b["w"] <= a["x"] + 1e-6
                or a["y"] + a["h"] <= b["y"] + 1e-6 or b["y"] + b["h"] <= a["y"] + 1e-6)


prs = Presentation(DECK)
print(f"{len(prs.slides)} slides · {prs.slide_width/EMU:.2f}in x {prs.slide_height/EMU:.2f}in\n")
problemas = 0

for i, slide in enumerate(prs.slides, 1):
    cs = caixas(slide)
    achados = []

    for c in cs:
        # fora do slide
        if c["x"] < -0.01 or c["y"] < -0.01 or c["x"] + c["w"] > LARG + 0.01 or c["y"] + c["h"] > ALT + 0.01:
            achados.append(f"FORA DO SLIDE: {c['texto'][:38]!r} em ({c['x']:.2f},{c['y']:.2f}) {c['w']:.2f}x{c['h']:.2f}")
        # estouro de texto
        if c["texto"].strip():
            alt, nl = altura_estimada(c)
            if alt > c["h"] + 0.06:
                achados.append(f"ESTOURO: {c['texto'][:38]!r} precisa ~{alt:.2f}in em caixa de {c['h']:.2f}in ({nl} linhas)")
        # margem inferior de conteudo textual
        if c["texto"].strip() and c["y"] + c["h"] > ALT - 0.25:
            achados.append(f"MARGEM: {c['texto'][:32]!r} termina em y={c['y']+c['h']:.2f} (slide {ALT})")

    # sobreposicao entre caixas com texto/forma (imagens de fundo ignoradas)
    conteudo = [c for c in cs if c["texto"].strip()]
    for a in range(len(conteudo)):
        for b in range(a + 1, len(conteudo)):
            if sobrepoe(conteudo[a], conteudo[b]):
                achados.append(f"SOBREPOSIÇÃO: {conteudo[a]['texto'][:24]!r} x {conteudo[b]['texto'][:24]!r}")

    if achados:
        problemas += len(achados)
        print(f"--- slide {i} ---")
        for a in achados:
            print("   ", a)

print(f"\n{'OK — nenhum problema geométrico' if not problemas else f'{problemas} pontos a revisar'}")
sys.exit(0)
